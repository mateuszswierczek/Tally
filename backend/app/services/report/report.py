from pptx import Presentation
from pptx.util import Inches
from app.services.analyzer.analyzer import Analyzer
from app.services.analyzer.schema import FrequencieTable, MAQTable, MatrixTable, Crosstable
import pandas as pd

class ReportGenerator:
    def __init__(self, analyzer:Analyzer) -> None:
        self.analyzer = analyzer
        self.pptx = Presentation()

    def generate_powerpoint(self):
        tables:list[Crosstable] = self.analyzer.tables
        for table in tables:
            slide = self.pptx.slides.add_slide(self.pptx.slide_layouts[0])
            try:
                rows, cols = table.percentage_table.shape
                rows += 1
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(3)
                pptx_table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                for col_idx, col_name in enumerate(table.percentage_table.columns):
                    pptx_table.cell(0, col_idx).text = str(col_name)         

                for row_idx in range(len(table.percentage_table)):
                    for col_idx in range(len(table.percentage_table.columns)):
                        pptx_table.cell(row_idx + 1, col_idx).text = str(table.percentage_table.iloc[row_idx, col_idx])        

            except:
                continue
        pass

df = pd.DataFrame({
    'Product': ['A', 'B', 'C'],
    'Sales': [100, 150, 200],
    'Profit': [20, 30, 45]
})

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

rows, cols = df.shape
rows += 1 

left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(3)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

for col_idx, col_name in enumerate(df.columns):
    table.cell(0, col_idx).text = str(col_name)

for row_idx in range(len(df)):
    for col_idx in range(len(df.columns)):
        table.cell(row_idx + 1, col_idx).text = str(df.iloc[row_idx, col_idx])

prs.save('output.pptx')